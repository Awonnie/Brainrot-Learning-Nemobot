from fastapi import FastAPI, Request, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi.responses import FileResponse, Response

from langchain.docstore.document import Document
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain.prompts import PromptTemplate
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains.summarize import load_summarize_chain
from langchain_community.llms import OpenAI

from fastapi.middleware.cors import CORSMiddleware


from pptx import Presentation
from elevenlabs.client import ElevenLabs
from pydub import AudioSegment
from dotenv import load_dotenv
import openai

import os
import json
import aiofiles
import requests
import uuid
import uvicorn
import re

# === Load Environment Variables ===
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

# === App Initialization ===
app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # or whitelist specific origins like ["https://your-frontend.com"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# === ElevenLabs Voice Settings ===
ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY")
client = ElevenLabs(api_key=ELEVENLABS_API_KEY)

VOICE_IDS = {
    "Blur Sotong": "ckdz71REaQKVx2gnOQjQ",
    "Smart Mugger": "aSXZu6bgEOS8MXVRzjPi"
}

# === Graceful favicon fallback ===
@app.get("/favicon.ico")
async def favicon():
    return Response(status_code=204)

# === Speech Synthesis ===
def synthesize_speech(text, speaker):
    voice_id = VOICE_IDS.get(speaker)
    if not voice_id:
        print(f"[ERROR] No voice ID found for speaker: {speaker}")
        return None

    try:
        audio_gen = client.text_to_speech.convert(
            voice_id=voice_id,
            model_id="eleven_multilingual_v2",
            text=text,
            output_format="mp3_44100_128"
        )

        os.makedirs("/tmp/audio/", exist_ok=True)
        audio_path = f"/tmp/audio/{uuid.uuid4()}.mp3"
        with open(audio_path, "wb") as f:
            f.write(b"".join(audio_gen))
        return audio_path
    except Exception as e:
        print(f"[ERROR] ElevenLabs synthesis failed: {e}")
        return None


# Temporary global state (in production, use session storage or DB)
duration_store = {}

@app.post("/set_duration")
async def set_duration(request: Request, duration: str = Form(...)):
    client_ip = request.client.host  # identify user (replace with real user/session tracking if needed)
    duration_store[client_ip] = duration
    return {"message": f"Duration set to {duration} for {client_ip}"}

# === Helper: Extract .pptx text ===
def extract_pptx_text(filepath):
    prs = Presentation(filepath)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def clean_temp_audio_dir():
    audio_dir = "/tmp/audio/"
    if os.path.exists(audio_dir):
        for filename in os.listdir(audio_dir):
            filepath = os.path.join(audio_dir, filename)
            try:
                os.remove(filepath)
            except Exception as e:
                print(f"[WARN] Failed to remove {filepath}: {e}")


# === Generate Brainrot Conversation ===
def generate_brainrot_convo(file_path, duration):
    print(f"[Brainrot] Processing file: {file_path}")

    if file_path.endswith(".pdf"):
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
        docs = text_splitter.split_documents(documents)

    elif file_path.endswith(".pptx"):
        raw_text = extract_pptx_text(file_path)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
        docs = text_splitter.create_documents([raw_text])

    else:
        raise ValueError("Unsupported file type. Please upload a .pdf or .pptx file.")

    convo_prompt_template = """
You are creating a fun and educational conversation between two students:
- Blur Sotong: A very confused student who doesn't understand anything.
- Smart Mugger: A hardworking student who explains concepts clearly and helpfully.

They are studying the following lecture notes together:
------------
{text}
------------

Generate a *natural-sounding*, back-and-forth conversation like this:

Blur Sotong: [confused question]
Smart Mugger: [clear explanation]
Blur Sotong: [follow-up confusion or curiosity]
Smart Mugger: [further clarification]

DO NOT include scene directions or markdown syntax like ** or ##.
Only produce lines that begin with either "Blur Sotong:" or "Smart Mugger:", followed by text.
End the conversation when you have explained all the key ideas.
Keep the duration of the conversation to {duration}
"""

    # Inject duration directly into the template string
    convo_prompt_with_duration = convo_prompt_template.replace("{duration}", duration)
    prompt = PromptTemplate(template=convo_prompt_with_duration, input_variables=["text"])

    conversation_data = []

    for i, chunk in enumerate(docs):
        print(f"[LLM] Processing chunk {i+1}/{len(docs)}")
        try:
            formatted_prompt = prompt.format(text=chunk.page_content)
            response = openai.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a helpful assistant generating a brainrot-style conversation between two students."},
                    {"role": "user", "content": formatted_prompt}
                ],
                model="gpt-4o"
            )
            result_text = response.choices[0].message.content.strip()
            lines = [line.strip() for line in result_text.split("\n") if ":" in line]

            for line in lines:
                speaker, message = line.split(":", 1)
                speaker = speaker.strip()
                message = message.strip()
                clean_speaker = re.sub(r"^[#\s*]+|[\s*]+$", "", speaker)
                audio_path = synthesize_speech(message, clean_speaker)
                conversation_data.append({
                    "speaker": clean_speaker,
                    "text": message,
                    "audio_path": audio_path
                })
        except Exception as e:
            print(f"[ERROR] Failed on chunk {i+1}: {e}")
            continue

    return conversation_data


@app.get("/")
async def home(request: Request):
    return templates.TemplateResponse("brainrot.html", {"request": request})

from fastapi.datastructures import FormData
from starlette.requests import Request

@app.post("/brainrot_chatbot")
async def brainrot_chatbot(
    request: Request,
    pdf_file: UploadFile = File(...),
):
    client_ip = request.client.host
    duration = duration_store.get(client_ip, "1 minute")

    docs_folder = "/tmp/docs/"
    os.makedirs(docs_folder, exist_ok=True)
    clean_temp_audio_dir()

    file_path = os.path.join(docs_folder, pdf_file.filename)
    file_bytes = await pdf_file.read()

    # Save the file to disk (so generate_text_conversation can use it)
    async with aiofiles.open(file_path, "wb") as f:
        await f.write(file_bytes)

    # Recreate a fresh UploadFile-like object with the same bytes for generate_summary
    from starlette.datastructures import UploadFile as StarletteUploadFile
    import io
    pdf_file_copy = StarletteUploadFile(
        filename=pdf_file.filename,
        file=io.BytesIO(file_bytes),
    )


    # Call generate_text_conversation
    text_conversation = generate_text_conversation(file_path, duration)

    # Save conversation
    os.makedirs("/tmp/output/", exist_ok=True)
    output_json = "/tmp/output/brainrot_convo.json"
    with open(output_json, "w", encoding="utf-8") as f:
        json.dump({"conversation": text_conversation}, f, indent=2)

    summary_json = await generate_summary(
        request=request,
        pdf_file=pdf_file_copy,
        summary_type="concise",
        audience="student"
    )



    return templates.TemplateResponse("brainrot.html", {
        "request": request,
        "conversation": text_conversation,
        "summary": summary_json.get("summary", "")
    })


def generate_text_conversation(file_path, duration):
    print(f"[Brainrot] Processing file: {file_path}")

    if file_path.endswith(".pdf"):
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
        docs = text_splitter.split_documents(documents)
    elif file_path.endswith(".pptx"):
        raw_text = extract_pptx_text(file_path)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
        docs = text_splitter.create_documents([raw_text])
    else:
        raise ValueError("Unsupported file type. Please upload a .pdf or .pptx file.")

    convo_prompt_template = """
You are creating a fun and educational conversation between two students:
- Blur Sotong: A very confused student who doesn't understand anything.
- Smart Mugger: A hardworking student who explains concepts clearly and helpfully.

They are studying the following lecture notes together:
------------
{text}
------------

Generate a *natural-sounding*, back-and-forth conversation like this:

Blur Sotong: [confused question]
Smart Mugger: [clear explanation]
Blur Sotong: [follow-up confusion or curiosity]
Smart Mugger: [further clarification]

DO NOT include scene directions or markdown syntax like ** or ##.
Only produce lines that begin with either "Blur Sotong:" or "Smart Mugger:", followed by text.
End the conversation when you have explained all the key ideas.
Keep the duration of the conversation to {duration}
"""

    # Inject duration directly into the template string
    convo_prompt_with_duration = convo_prompt_template.replace("{duration}", duration)
    prompt = PromptTemplate(template=convo_prompt_with_duration, input_variables=["text"])

    conversation_data = []

    for i, chunk in enumerate(docs):
        print(f"[LLM] Processing chunk {i+1}/{len(docs)}")
        try:
            formatted_prompt = prompt.format(text=chunk.page_content)
            response = openai.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a helpful assistant generating a brainrot-style conversation between two students."},
                    {"role": "user", "content": formatted_prompt}
                ],
                model="gpt-4o"
            )
            result_text = response.choices[0].message.content.strip()
            lines = [line.strip() for line in result_text.split("\n") if ":" in line]

            for line in lines:
                speaker, message = line.split(":", 1)
                speaker = speaker.strip()
                message = message.strip()
                clean_speaker = re.sub(r"^[#\s*]+|[\s*]+$", "", speaker)
                conversation_data.append({
                    "speaker": clean_speaker,
                    "text": message
                    # No audio_path here yet
                })
        except Exception as e:
            print(f"[ERROR] Failed on chunk {i+1}: {e}")
            continue

    return conversation_data

@app.post("/generate_full_audio")
async def generate_full_audio(request: Request):
    print("[🔊] Generating audio from existing conversation...")
    
    json_path = "/tmp/output/brainrot_convo.json"
    if not os.path.exists(json_path):
        return {"status": "error", "message": "No conversation found. Please upload a document first."}

    with open(json_path, "r", encoding="utf-8") as f:
        convo_data = json.load(f)

    conversation = convo_data.get("conversation", [])
    podcast_audio = AudioSegment.empty()

    for idx, entry in enumerate(conversation):
        # Generate audio for each line
        audio_path = synthesize_speech(entry["text"], entry["speaker"])
        if audio_path:
            conversation[idx]["audio_path"] = audio_path
            segment = AudioSegment.from_mp3(audio_path)
            podcast_audio += segment + AudioSegment.silent(duration=300)

    # Save updated conversation with audio paths
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({"conversation": conversation}, f, indent=2)

    final_audio_path = "/tmp/brainrot_full_audio.mp3"
    podcast_audio.export(final_audio_path, format="mp3")

    return {
        "status": "success",
        "audio_file": "/download_audio"
    }

@app.get("/download_audio")
async def download_audio():
    return FileResponse("/tmp/brainrot_full_audio.mp3", media_type="audio/mpeg", filename="brainrot_full_audio.mp3")

@app.get("/get_audio_link")
async def get_audio_link(request: Request):
    client_ip = request.client.host
    audio_file_path = "/tmp/brainrot_full_audio.mp3"

    if os.path.exists(audio_file_path):
        return {
            "status": "success",
            "message": "Audio generated successfully!",
            "audio_url": "https://brainrot-learning-4052.uc.r.appspot.com/download_audio"
        }
    else:
        return {
            "status": "error",
            "message": "Audio file not found. Please try again after processing the PDF."
        }



@app.post("/generate_summary")
async def generate_summary(
    request: Request, 
    pdf_file: UploadFile = File(...),
    summary_type: str = Form("concise"),
    audience: str = Form("student")
):
    print("[Brainrot] Received summary generation request")

    docs_folder = "/tmp/docs/"
    output_folder = "/tmp/output/"
    os.makedirs(docs_folder, exist_ok=True)
    os.makedirs(output_folder, exist_ok=True)

    file_path = os.path.join(docs_folder, pdf_file.filename)
    output_json = os.path.join(output_folder, "brainrot_summary.json")

    # ✅ Save the uploaded file to disk first
    async with aiofiles.open(file_path, "wb") as f:
        content = await pdf_file.read()
        if not content:
            print("[ERROR] Empty file uploaded.")
            return {"error": "Uploaded file is empty"}
        await f.write(content)

    if not file_path.endswith(".pdf"):
        return {"error": "Only PDF files are supported for summarization"}

    try:
        # ✅ Now the file has content, and we can load it
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
        docs = text_splitter.split_documents(documents)

        summary_style = (
            "a comprehensive, point-by-point summary"
            if summary_type == "detailed"
            else "a brief 3-5 sentence summary"
        )
        audience_desc = (
            "for a professional audience using formal language"
            if audience == "professional"
            else "for a student audience using simple language"
        )

        prompt_template = f"""
You are a helpful assistant who writes summaries of academic documents.

Write {summary_style} of the following text {audience_desc}:
------------
{{text}}
------------
SUMMARY:
"""
        prompt = PromptTemplate(template=prompt_template, input_variables=["text"])
        summary_chunks = []

        for i, chunk in enumerate(docs):
            print(f"[LLM] Summarizing chunk {i+1}/{len(docs)}")
            try:
                formatted_prompt = prompt.format(text=chunk.page_content)
                response = openai.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant that writes summaries based on provided academic texts."},
                        {"role": "user", "content": formatted_prompt}
                    ]
                )
                result_text = response.choices[0].message.content.strip()
                summary_chunks.append(result_text)
            except Exception as e:
                print(f"[ERROR] Failed on chunk {i+1}: {e}")
                continue

        final_summary = "\n\n".join(summary_chunks)

        with open(output_json, "w", encoding="utf-8") as f:
            json.dump({"summary": final_summary}, f, indent=2)

        return {
            "message": "Summary generated successfully",
            "summary": final_summary,
            "summary_path": "/download_summary"
        }

    except Exception as e:
        print(f"[ERROR] Summary generation failed: {e}")
        return {"error": "Failed to generate summary"}


    
@app.get("/download_summary")
async def download_summary():
    summary_file_path = "/tmp/output/brainrot_summary.json"
    
    if not os.path.exists(summary_file_path):
        return {"summary": ""}  # still valid JSON, but no summary

    try:
        with open(summary_file_path, "r") as f:
            summary_data = json.load(f)
        return summary_data
    except Exception as e:
        print(f"[ERROR] Failed to load summary JSON: {e}")
        return {"summary": ""}  # ensure client always gets valid JSON


if __name__ == "__main__":
    print("[Startup] Launching FastAPI server...")
    uvicorn.run("app:app", host="0.0.0.0", port=int(os.environ.get("PORT", 8080)))